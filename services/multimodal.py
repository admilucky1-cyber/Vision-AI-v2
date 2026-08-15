"""
Vision AI v2.0 - Multi-Modal File Processor (multimodal.py)
============================================================
Unified entry point for processing all supported file types.
Handles: PDF (text + OCR), Images (captioning), Text files, Office docs,
Audio (Whisper paid/free), URLs.

Supported formats:
- PDF: pdfplumber -> PyPDF2 -> pdfminer -> OCR (pytesseract)
- Images: JPG, PNG, BMP, WEBP, TIFF (BLIP captioning + OCR fallback)
- Office: DOCX, PPTX, XLSX
- Text: TXT, MD, CSV, JSON, XML, HTML, LOG, PY
- Audio: MP3, WAV, M4A, OGG (Whisper paid/free/local)
- URLs: HTTP/HTTPS (BeautifulSoup scraping)
"""

import os
import io
import tempfile
import traceback
import logging
import re
import base64
from pathlib import Path
from typing import Optional, Dict, Any, List
import requests
from dotenv import load_dotenv

load_dotenv()

# ==========================================================
# LOGGING SETUP
# ==========================================================
logger = logging.getLogger("vision-ai.multimodal")

# ==========================================================
# CONFIGURATION
# ==========================================================
HF_TOKEN = os.getenv("HF_TOKEN") or os.getenv("HUGGINGFACE_API_KEY")
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY") or os.getenv("GEMINI_API_KEY")
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")

# Image caption / vision endpoints (tried in order)
HF_IMAGE_MODELS = [
    "https://api-inference.huggingface.co/models/Salesforce/blip-image-captioning-large",
    "https://api-inference.huggingface.co/models/Salesforce/blip-image-captioning-base",
    "https://router.huggingface.co/hf-inference/models/Salesforce/blip-image-captioning-base",
]

# Supported file extensions
IMAGE_EXTENSIONS = {'jpg', 'jpeg', 'png', 'bmp', 'webp', 'tiff', 'tif'}
TEXT_EXTENSIONS = {'txt', 'md', 'csv', 'json', 'xml', 'html', 'log', 'py'}
OFFICE_EXTENSIONS = {'docx', 'doc', 'pptx', 'ppt', 'xlsx', 'xls'}
AUDIO_EXTENSIONS = {'mp3', 'wav', 'm4a', 'ogg', 'flac'}
URL_PATTERN = re.compile(r'^https?://')

# Max file sizes
MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB
MAX_TEXT_LENGTH = 100000  # 100K characters

# ==========================================================
# EXCEPTIONS
# ==========================================================
class FileProcessingError(Exception):
    """Custom exception for file processing failures."""
    pass

class FileTooLargeError(FileProcessingError):
    """Raised when file exceeds size limit."""
    pass

class UnsupportedFileTypeError(FileProcessingError):
    """Raised when file type is not supported."""
    pass

# ==========================================================
# MAIN PROCESSOR (✅ FIXED: SYNCHRONOUS)
# ==========================================================
def process_uploaded_file(file) -> str:
    """
    Process any uploaded file and return text representation.

    Args:
        file: FastAPI UploadFile object

    Returns:
        Extracted text content or error message
    """
    try:
        # Check file size
        if file.size and file.size > MAX_FILE_SIZE:
            raise FileTooLargeError(f"File '{file.filename}' exceeds {MAX_FILE_SIZE//1024//1024}MB limit")

        # ✅ FIX: Use synchronous file read
        content = file.file.read()

        if not content:
            logger.warning(f"File '{file.filename}' is empty")
            return f"[File '{file.filename}' is empty.]"

        filename = file.filename or "unknown"
        ext = Path(filename).suffix.lower().lstrip('.')

        logger.info(f"Processing file: {filename} (size: {len(content)} bytes, type: {ext})")

        # Route to appropriate handler
        if ext == 'pdf':
            return _process_pdf(content, filename)
        elif ext in IMAGE_EXTENSIONS:
            return _process_image(content, filename)
        elif ext in TEXT_EXTENSIONS:
            return _process_text(content, filename)
        elif ext in ['docx', 'doc']:
            return _process_docx(content, filename)
        elif ext in ['pptx', 'ppt']:
            return _process_pptx(content, filename)
        elif ext in ['xlsx', 'xls']:
            return _process_excel(content, filename)
        elif ext in AUDIO_EXTENSIONS:
            return _process_audio(content, filename)
        elif re.match(URL_PATTERN, filename):
            # If filename is actually a URL
            return _process_url(content, filename)
        else:
            logger.warning(f"Unsupported file type: .{ext}")
            return f"[Unsupported file type: .{ext}]"

    except FileTooLargeError as e:
        logger.error(f"File too large: {file.filename}")
        return str(e)
    except Exception as e:
        logger.error(f"Error processing {file.filename}: {e}")
        traceback.print_exc()
        return f"[Could not process '{file.filename}': {str(e)}]"

# ==========================================================
# PDF PROCESSING
# ==========================================================
def _pdf_text_quality(text: str, page_count_hint: int = 0) -> dict:
    """Heuristic: cover-only extracts score low even if they have >20 chars."""
    raw = (text or "").strip()
    n = len(raw)
    pages = max(page_count_hint, raw.count("--- Page "), 1)
    avg = n / pages
    low = raw.lower()
    front_only = 0
    for marker in (
        "candidate surname", "centre number", "pearson edexcel", "turn over",
        "total marks", "instructions", "information", "advice",
        "equation", "formulae", "constants", "data sheet",
    ):
        if marker in low:
            front_only += 1
    has_questions = bool(re.search(r"(?i)\b(q\s*\d+|question\s*\d+|\d+\s*\([a-d]\))", raw))
    score = 0
    if n > 2000:
        score += 2
    if avg > 400:
        score += 2
    if has_questions:
        score += 3
    if front_only >= 3 and not has_questions:
        score -= 3
    if n < 800:
        score -= 2
    return {"score": score, "chars": n, "pages": pages, "has_questions": has_questions, "avg": avg}


def _process_pdf(content: bytes, filename: str) -> str:
    """Process PDF with quality gate so exam papers are not reduced to cover sheets."""
    tmp_path = _save_temp_file(content, suffix=".pdf")

    try:
        extracted_text = _extract_text_from_pdf(tmp_path)
        page_hint = 0
        try:
            import fitz
            page_hint = len(fitz.open(tmp_path))
        except Exception:
            try:
                from PyPDF2 import PdfReader
                page_hint = len(PdfReader(tmp_path).pages)
            except Exception:
                page_hint = 0

        q = _pdf_text_quality(extracted_text or "", page_hint)
        logger.info(
            "PDF extract quality %s: score=%s chars=%s pages~%s questions=%s",
            filename, q["score"], q["chars"], q["pages"], q["has_questions"],
        )

        if extracted_text and q["score"] >= 2 and q["chars"] > 100:
            logger.info(f"PDF text extracted: {filename} ({q['chars']} chars)")
            return f"[PDF Content - {filename}]\n\n{extracted_text}"

        ocr_text = _ocr_pdf(tmp_path, filename)
        if ocr_text and len(ocr_text.strip()) > len((extracted_text or "").strip()):
            logger.info(f"PDF OCR preferred: {filename} ({len(ocr_text)} chars)")
            return f"[PDF Content (OCR) - {filename}]\n\n{ocr_text}"

        if extracted_text and len(extracted_text.strip()) > 20:
            warn = ""
            if q["score"] < 2:
                warn = (
                    "\n[WARNING: Text extract may be incomplete (cover/equations only). "
                    "If questions are missing, the PDF may be image-based — ensure OCR is installed.]\n"
                )
            logger.warning(f"PDF weak extract used: {filename} ({q['chars']} chars, score={q['score']})")
            return f"[PDF Content - {filename}]{warn}\n\n{extracted_text}"

        if ocr_text and len(ocr_text.strip()) > 20:
            logger.info(f"PDF OCR processed: {filename} ({len(ocr_text)} chars)")
            return f"[PDF Content (OCR) - {filename}]\n\n{ocr_text}"

        return (
            f"[PDF: {filename}] Could not extract text. "
            "Likely a scanned PDF. Install tesseract-ocr + poppler-utils, or paste questions as text. "
            "On Docker this is included; on Windows install Tesseract and set PATH."
        )
    finally:
        _cleanup_temp_file(tmp_path)



def _extract_text_from_pdf(pdf_path: str) -> Optional[str]:
    """Extract text using multiple PDF libraries."""
    methods = [_try_pymupdf, _try_pdfplumber, _try_pypdf2, _try_pdfminer]

    best = ""
    for method in methods:
        try:
            text = method(pdf_path)
            if text and len(text.strip()) > len(best):
                best = text.strip()
                logger.debug(f"PDF method {method.__name__}: {len(best)} chars")
        except Exception as e:
            logger.debug(f"PDF extraction method failed: {e}")
            continue
    return best or None


def _try_pymupdf(pdf_path: str) -> str:
    """PyMuPDF often recovers more text from exam PDFs than pdfplumber alone."""
    try:
        import fitz
        doc = fitz.open(pdf_path)
        pages_text = []
        for i, page in enumerate(doc, 1):
            text = page.get_text("text") or ""
            if text.strip():
                pages_text.append(f"--- Page {i} ---\n{text}")
        doc.close()
        return "\n\n".join(pages_text)
    except ImportError:
        raise
    except Exception as e:
        logger.debug(f"pymupdf failed: {e}")
        raise


def _try_pdfplumber(pdf_path: str) -> str:
    try:
        import pdfplumber
        with pdfplumber.open(pdf_path) as pdf:
            pages_text = []
            for i, page in enumerate(pdf.pages, 1):
                text = page.extract_text()
                if text and text.strip():
                    pages_text.append(f"--- Page {i} ---\n{text}")
            return "\n\n".join(pages_text)
    except ImportError:
        raise
    except Exception as e:
        logger.debug(f"pdfplumber failed: {e}")
        raise

def _try_pypdf2(pdf_path: str) -> str:
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(pdf_path)
        pages_text = []
        for i, page in enumerate(reader.pages, 1):
            text = page.extract_text()
            if text and text.strip():
                pages_text.append(f"--- Page {i} ---\n{text}")
        return "\n\n".join(pages_text)
    except ImportError:
        raise
    except Exception as e:
        logger.debug(f"PyPDF2 failed: {e}")
        raise

def _try_pdfminer(pdf_path: str) -> str:
    try:
        from pdfminer.high_level import extract_text
        return extract_text(pdf_path)
    except ImportError:
        raise
    except Exception as e:
        logger.debug(f"pdfminer failed: {e}")
        raise

def _ocr_pdf(pdf_path: str, filename: str) -> Optional[str]:
    """Rasterize PDF pages and OCR (needs poppler + tesseract)."""
    try:
        from pdf2image import convert_from_path
        import pytesseract
    except ImportError as e:
        logger.warning(f"OCR deps missing: {e}")
        return None
    try:
        # Limit pages for speed/memory; high DPI for exam papers
        images = convert_from_path(pdf_path, dpi=220, first_page=1, last_page=int(__import__('os').getenv('PDF_OCR_MAX_PAGES', '40')))
        chunks = []
        for i, img in enumerate(images, 1):
            try:
                text = pytesseract.image_to_string(img) or ""
                text = text.strip()
                if text:
                    chunks.append(f"--- Page {i} ---\n{text}")
            except Exception as pe:
                logger.debug(f"OCR page {i} failed: {pe}")
        if chunks:
            return "\n\n".join(chunks)
    except Exception as e:
        logger.warning(f"PDF OCR failed for {filename}: {e}")
    return None



def _format_pdf_error(filename: str) -> str:
    return (
        f"[Could not read '{filename}'. This PDF may be:]\n"
        f"• A scanned document (images only, no text)\n"
        f"• Password protected\n"
        f"• Corrupted or malformed\n\n"
        f"Solutions:\n"
        f"1. If scanned: Install OCR support (pip install pytesseract pdf2image)\n"
        f"2. If protected: Remove password before uploading\n"
        f"3. Copy-paste the text directly into chat"
    )

# ==========================================================
# IMAGE PROCESSING
# ==========================================================

def _process_image(content: bytes, filename: str) -> str:
    """
    Analyze image with multi-provider fallback:
    1) Local OCR (pytesseract) — free, offline
    2) Gemini vision — if GOOGLE_API_KEY set
    3) Hugging Face BLIP caption — if HF_TOKEN set
    Always returns usable context for the chat LLM.
    """
    parts: List[str] = [f"[IMAGE: {filename}]"]
    ocr_text = _ocr_image_bytes(content)
    if ocr_text:
        parts.append("### Text visible in image (OCR)\n" + ocr_text[:8000])

    caption = None
    # Gemini vision
    if GOOGLE_API_KEY:
        try:
            caption = _gemini_describe_image(content, filename)
            if caption:
                parts.append("### Visual description (Gemini)\n" + caption[:6000])
        except Exception as e:
            logger.warning(f"Gemini image analysis failed: {e}")

    # HF BLIP
    if not caption and HF_TOKEN:
        try:
            caption = _hf_caption_image(content)
            if caption:
                parts.append("### Caption (Hugging Face)\n" + caption[:4000])
        except Exception as e:
            logger.warning(f"HF caption failed: {e}")

    if len(parts) == 1:
        # Still give the model something to work with
        parts.append(
            "Could not run OCR or cloud vision (install tesseract-ocr and/or set "
            "GOOGLE_API_KEY or HF_TOKEN). File was received; ask the user to describe it."
        )
        return "\n\n".join(parts)

    parts.append(
        "Use the OCR text and description above to answer the user's question about this image."
    )
    return "\n\n".join(parts)


def _ocr_image_bytes(content: bytes) -> Optional[str]:
    """Local OCR via pytesseract + Pillow (no API key)."""
    try:
        from PIL import Image
        import pytesseract
        img = Image.open(io.BytesIO(content))
        if img.mode not in ("RGB", "L"):
            img = img.convert("RGB")
        # Slight upscale helps small exam scans
        w, h = img.size
        if max(w, h) < 1200:
            scale = 1200 / max(w, h)
            img = img.resize((int(w * scale), int(h * scale)), Image.Resampling.LANCZOS)
        text = pytesseract.image_to_string(img) or ""
        text = text.strip()
        if len(text) >= 3:
            logger.info(f"OCR extracted {len(text)} chars from image")
            return text
    except ImportError:
        logger.warning("Pillow/pytesseract not available for image OCR")
    except Exception as e:
        logger.warning(f"Image OCR failed: {e}")
    return None


def _gemini_describe_image(content: bytes, filename: str) -> Optional[str]:
    """Use Google Gemini multimodal (free tier with API key)."""
    import base64
    b64 = base64.b64encode(content).decode("ascii")
    # Detect mime
    mime = "image/jpeg"
    low = (filename or "").lower()
    if low.endswith(".png"):
        mime = "image/png"
    elif low.endswith(".webp"):
        mime = "image/webp"
    elif low.endswith(".gif"):
        mime = "image/gif"

    models = [
        "gemini-2.0-flash",
        "gemini-2.5-flash",
        "gemini-1.5-flash",
        "gemini-1.5-flash-latest",
    ]
    prompt = (
        "You are a careful visual analyst. Describe EVERYTHING visible in this image with high detail.\n"
        "- If screenshot: name the apps/windows, UI text, colors, layout.\n"
        "- If logo/branding: describe symbols, colors, slogan text exactly.\n"
        "- If exam/diagram: extract all numbers, labels, options A-D, axes, equations.\n"
        "- If photo: subjects, setting, text on signs.\n"
        "Never give a vague summary. Prefer specific observed details over guesses."
    )
    for model in models:
        try:
            url = (
                f"https://generativelanguage.googleapis.com/v1beta/models/"
                f"{model}:generateContent?key={GOOGLE_API_KEY}"
            )
            body = {
                "contents": [{
                    "parts": [
                        {"text": prompt},
                        {"inline_data": {"mime_type": mime, "data": b64}},
                    ]
                }],
                "generationConfig": {"temperature": 0.2, "maxOutputTokens": 2048},
            }
            resp = requests.post(url, json=body, timeout=60)
            if resp.status_code != 200:
                logger.debug(f"Gemini {model} image: {resp.status_code} {resp.text[:200]}")
                continue
            data = resp.json()
            cands = data.get("candidates") or []
            if not cands:
                continue
            parts = (cands[0].get("content") or {}).get("parts") or []
            text = "".join(p.get("text", "") for p in parts).strip()
            if text:
                logger.info(f"Gemini vision OK via {model} ({len(text)} chars)")
                return text
        except Exception as e:
            logger.debug(f"Gemini {model} error: {e}")
            continue
    return None


def _hf_caption_image(content: bytes) -> Optional[str]:
    """Hugging Face BLIP captioning."""
    headers = {"Authorization": f"Bearer {HF_TOKEN}"}
    for url in HF_IMAGE_MODELS:
        try:
            resp = requests.post(url, headers=headers, data=content, timeout=45)
            if resp.status_code == 503:
                # model loading
                import time
                time.sleep(3)
                resp = requests.post(url, headers=headers, data=content, timeout=60)
            if resp.status_code != 200:
                logger.debug(f"HF caption {resp.status_code}: {resp.text[:150]}")
                continue
            data = resp.json()
            if isinstance(data, list) and data:
                text = data[0].get("generated_text") or data[0].get("caption") or ""
            elif isinstance(data, dict):
                text = data.get("generated_text") or data.get("caption") or str(data)
            else:
                text = str(data)
            text = (text or "").strip()
            if text:
                return text
        except Exception as e:
            logger.debug(f"HF caption error: {e}")
            continue
    return None



def _process_audio(content: bytes, filename: str) -> str:
    """
    Transcribe audio using Whisper.
    Priority:
    1. OpenAI API (paid, best quality)
    2. Hugging Face API (free, good quality)
    3. Local Whisper (free, offline)
    """
    tmp_path = _save_temp_file(content, suffix='.mp3')

    try:
        # Option 1: OpenAI Whisper (paid)
        if OPENAI_API_KEY:
            try:
                logger.info(f"Using OpenAI Whisper (paid) for {filename}")
                import openai
                openai.api_key = OPENAI_API_KEY
                with open(tmp_path, "rb") as f:
                    response = openai.Audio.transcribe(
                        model="whisper-1",
                        file=f
                    )
                text = response.get("text", "").strip()
                if text:
                    logger.info(f"OpenAI Whisper transcribed: {filename} ({len(text)} chars)")
                    return f"[Audio Transcription (OpenAI): {filename}]\n\n{text}"
            except Exception as e:
                logger.warning(f"OpenAI Whisper failed: {e}. Falling back to HF/local.")

        # Option 2: Hugging Face Whisper (free)
        if HF_TOKEN:
            try:
                logger.info(f"Using Hugging Face Whisper (free) for {filename}")
                import requests
                headers = {"Authorization": f"Bearer {HF_TOKEN}"}
                with open(tmp_path, "rb") as f:
                    resp = requests.post(
                        "https://api-inference.huggingface.co/models/openai/whisper-large-v3",
                        headers=headers,
                        data=f.read(),
                        timeout=60
                    )
                if resp.status_code == 200:
                    result = resp.json()
                    text = result.get("text", "").strip()
                    if text:
                        logger.info(f"HF Whisper transcribed: {filename} ({len(text)} chars)")
                        return f"[Audio Transcription (Hugging Face): {filename}]\n\n{text}"
                else:
                    logger.warning(f"HF Whisper returned {resp.status_code}. Falling back to local.")
            except Exception as e:
                logger.warning(f"HF Whisper failed: {e}. Falling back to local.")

        # Option 3: Local Whisper (free, offline)
        try:
            logger.info(f"Using local Whisper (free, offline) for {filename}")
            import whisper
            model = whisper.load_model("base")  # Use "tiny" for speed, "large" for accuracy
            result = model.transcribe(tmp_path)
            text = result.get("text", "").strip()
            if text:
                logger.info(f"Local Whisper transcribed: {filename} ({len(text)} chars)")
                return f"[Audio Transcription (Local): {filename}]\n\n{text}"
        except ImportError:
            logger.warning("Local Whisper not installed. Skipping.")
        except Exception as e:
            logger.warning(f"Local Whisper failed: {e}")

        # If all fail
        return f"[Could not transcribe audio '{filename}'. No Whisper model available.]"

    finally:
        _cleanup_temp_file(tmp_path)

# ==========================================================
# URL PROCESSING
# ==========================================================
def _process_url(content: bytes, filename: str) -> str:
    """Fetch and extract text from a URL."""
    try:
        from bs4 import BeautifulSoup
        url = content.decode('utf-8').strip()
        if not url.startswith(('http://', 'https://')):
            return f"[Invalid URL: {url}]"
        
        headers = {"User-Agent": "VisionAI/2.0"}
        resp = requests.get(url, headers=headers, timeout=15)
        soup = BeautifulSoup(resp.text, 'html.parser')
        
        # Remove script and style tags
        for script in soup(["script", "style"]):
            script.decompose()
        
        title = soup.title.string if soup.title else "No title"
        text = soup.get_text()[:10000]  # Limit to 10K chars
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        text = '\n'.join(chunk for chunk in chunks if chunk)
        
        return f"[URL: {url}]\nTitle: {title}\n\n{text}"
    except ImportError:
        return "[URL fetching requires: pip install beautifulsoup4]"
    except Exception as e:
        return f"[Error fetching URL: {str(e)}]"

# ==========================================================
# OFFICE DOCUMENTS
# ==========================================================
def _process_docx(content: bytes, filename: str) -> str:
    try:
        from docx import Document
        doc = Document(io.BytesIO(content))
        texts = [para.text for para in doc.paragraphs if para.text.strip()]
        logger.info(f"DOCX processed: {filename} ({len(texts)} paragraphs)")
        return f"[Word Document: {filename}]\n\n" + "\n".join(texts)
    except ImportError:
        logger.warning("python-docx not installed")
        return "[Word document support requires: pip install python-docx]"
    except Exception as e:
        logger.error(f"DOCX extraction error: {e}")
        return f"[Error extracting DOCX: {str(e)}]"

def _process_pptx(content: bytes, filename: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(content))
        texts = []
        max_slides = int(os.getenv("PPTX_MAX_SLIDES", "80") or 80)
        total = len(prs.slides)
        for slide_num, slide in enumerate(prs.slides, 1):
            if slide_num > max_slides:
                texts.append(f"[… {total - max_slides} more slides omitted — ask for a slide range]")
                break
            slide_text = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            if slide_text:
                texts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_text))
        logger.info(f"PPTX processed: {filename} ({min(total, max_slides)}/{total} slides)")
        return f"[PowerPoint: {filename} — {total} slides]\n\n" + "\n\n".join(texts)
    except ImportError:
        logger.warning("python-pptx not installed")
        return "[PowerPoint support requires: pip install python-pptx]"
    except Exception as e:
        logger.error(f"PPTX extraction error: {e}")
        return f"[Error extracting PPTX: {str(e)}]"

def _process_excel(content: bytes, filename: str) -> str:
    """
    Large workbooks: sample head rows per sheet + summary stats so millions of
    cells do not explode memory / LLM context. Full dump only for small sheets.
    """
    MAX_ROWS_PER_SHEET = int(os.getenv("EXCEL_MAX_ROWS", "200") or 200)
    MAX_SHEETS = int(os.getenv("EXCEL_MAX_SHEETS", "20") or 20)
    MAX_CHARS = int(os.getenv("EXCEL_MAX_CHARS", "80000") or 80000)
    try:
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(content), data_only=True, read_only=True)
        texts = []
        total_rows_seen = 0
        for sheet_i, sheet in enumerate(wb.worksheets):
            if sheet_i >= MAX_SHEETS:
                texts.append(f"[… {len(wb.worksheets) - MAX_SHEETS} more sheets omitted]")
                break
            sheet_text = []
            row_count = 0
            shown = 0
            for row in sheet.iter_rows(values_only=True):
                row_count += 1
                total_rows_seen += 1
                if shown >= MAX_ROWS_PER_SHEET:
                    continue
                cells = [str(cell) for cell in row if cell is not None]
                if not cells:
                    continue
                sheet_text.append(" | ".join(cells))
                shown += 1
            header = f"[Sheet: {sheet.title} — {row_count} rows"
            if row_count > MAX_ROWS_PER_SHEET:
                header += f", showing first {MAX_ROWS_PER_SHEET}"
            header += "]"
            body = "\n".join(sheet_text) if sheet_text else "(empty)"
            texts.append(f"{header}\n{body}")
        try:
            wb.close()
        except Exception:
            pass
        out = f"[Excel: {filename} | ~{total_rows_seen} data rows scanned]\n\n" + "\n\n".join(texts)
        if len(out) > MAX_CHARS:
            out = out[:MAX_CHARS] + "\n\n[Excel truncated — ask for a specific sheet/column for more detail]"
        logger.info(f"Excel processed: {filename} (rows~{total_rows_seen})")
        return out
    except ImportError:
        logger.warning("openpyxl not installed")
        return "[Excel support requires: pip install openpyxl]"
    except Exception as e:
        logger.error(f"Excel extraction error: {e}")
        return f"[Error extracting Excel: {str(e)}]"

# ==========================================================
# TEXT PROCESSING
# ==========================================================
def _process_text(content: bytes, filename: str) -> str:
    encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
    for encoding in encodings:
        try:
            text = content.decode(encoding)
            # Large CSV: sample head + tail + line count (do not load millions of lines into LLM)
            lower = (filename or "").lower()
            if lower.endswith(".csv") or lower.endswith(".tsv"):
                lines = text.splitlines()
                n = len(lines)
                max_lines = int(os.getenv("CSV_MAX_LINES", "300") or 300)
                if n > max_lines:
                    head = lines[: max_lines // 2]
                    tail = lines[-(max_lines // 2) :]
                    text = (
                        f"[CSV: {filename} — {n} lines total; showing first/last samples]\n"
                        + "\n".join(head)
                        + f"\n\n… ({n - max_lines} lines omitted) …\n\n"
                        + "\n".join(tail)
                    )
                else:
                    text = f"[CSV: {filename} — {n} lines]\n" + text
            if len(text) > MAX_TEXT_LENGTH:
                text = text[:MAX_TEXT_LENGTH] + "\n\n[Text truncated due to length]"
            logger.info(f"Text file processed: {filename} ({len(text)} chars)")
            return f"[File: {filename}]\n\n{text}"
        except UnicodeDecodeError:
            continue
    logger.warning(f"Could not decode text file: {filename}")
    return f"[Could not decode text file '{filename}'. Unsupported encoding.]"

# ==========================================================
# UTILITY FUNCTIONS
# ==========================================================
def _save_temp_file(content: bytes, suffix: str = '') -> str:
    """Save content to a temporary file and return path."""
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(content)
        tmp_path = tmp.name
        logger.debug(f"Temporary file created: {tmp_path}")
        return tmp_path

def _cleanup_temp_file(path: str) -> None:
    """Remove temporary file if it exists."""
    try:
        if os.path.exists(path):
            os.unlink(path)
            logger.debug(f"Temporary file removed: {path}")
    except Exception as e:
        logger.warning(f"Failed to remove temporary file {path}: {e}")

def check_file_size(content: bytes, max_size: int = MAX_FILE_SIZE) -> bool:
    """Check if file size is within limits."""
    return len(content) <= max_size

def get_file_extension(filename: str) -> str:
    """Get file extension from filename."""
    return Path(filename).suffix.lower().lstrip('.')

# ==========================================================
# EXPORTS
# ==========================================================
__all__ = [
    "process_uploaded_file",
    "FileProcessingError",
    "FileTooLargeError",
    "UnsupportedFileTypeError",
    "check_file_size",
    "get_file_extension",
]

logger.info("👁️ Vision AI Multi-Modal Processor v2.0 - Ready")